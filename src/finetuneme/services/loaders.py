"""
Universal document loaders for different file types.
Each loader outputs DocumentChunk objects compatible with the existing chunking logic.
"""
from abc import ABC, abstractmethod
from typing import List, Dict
from pathlib import Path
import re
import io
import base64

# Import for specific loaders
import fitz  # PyMuPDF for PDFs
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    from PIL import Image
except ImportError:
    Image = None

from finetuneme.services.ingestion import DocumentChunk, clean_text, chunk_text_semantic


class DocumentLoader(ABC):
    """Abstract base class for document loaders"""

    @abstractmethod
    def load(self, file_path: str) -> List[DocumentChunk]:
        """
        Load and parse a document into chunks.

        Args:
            file_path: Path to the document file

        Returns:
            List of DocumentChunk objects
        """
        pass

    @staticmethod
    def supports(file_extension: str) -> bool:
        """Check if this loader supports the given file extension"""
        return False


class PDFLoader(DocumentLoader):
    """Loader for PDF files using PyMuPDF"""

    SUPPORTED_EXTENSIONS = {'.pdf'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in PDFLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load PDF file and extract text chunks, with image extraction for scanned pages"""
        doc = fitz.open(file_path)
        all_chunks = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            cleaned_text = clean_text(text)

            # Check if this is likely a scanned page (very little text)
            is_scanned = len(cleaned_text) < 50
            page_images = None

            if is_scanned:
                # Render page as image
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x resolution for better quality
                    img_bytes = pix.tobytes("jpeg")
                    img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                    page_images = [img_base64]

                    # Update text to indicate it's an image
                    if not cleaned_text:
                        cleaned_text = "[Scanned Page Image]"
                except Exception as e:
                    print(f"Warning: Could not extract image from page {page_num + 1}: {str(e)}")

            if not cleaned_text and not page_images:
                continue

            # Chunk the page text
            if cleaned_text:
                page_chunks = chunk_text_semantic(cleaned_text)
            else:
                page_chunks = ["[Scanned Page Image]"]

            # Create DocumentChunk objects
            for chunk_text in page_chunks:
                chunk = DocumentChunk(
                    text=chunk_text,
                    page_num=page_num + 1,
                    metadata={
                        "source": file_path,
                        "total_pages": len(doc),
                        "file_type": "pdf",
                        "is_scanned": is_scanned
                    },
                    images=page_images
                )
                all_chunks.append(chunk)

        doc.close()
        return all_chunks


class WordLoader(DocumentLoader):
    """Loader for Microsoft Word documents (.docx)"""

    SUPPORTED_EXTENSIONS = {'.docx', '.doc'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in WordLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load Word document and extract text chunks"""
        if DocxDocument is None:
            raise ImportError("python-docx is required for Word document support. Install with: pip install python-docx")

        if file_path.lower().endswith('.doc'):
             raise ValueError("Legacy .doc files (Word 97-2003) are not supported by python-docx. Please convert to modern .docx format.")

        doc = DocxDocument(file_path)
        all_chunks = []

        # Extract text from all paragraphs
        full_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    full_text.append(" | ".join(row_text))

        combined_text = "\n\n".join(full_text)
        cleaned_text = clean_text(combined_text)

        if not cleaned_text:
            return []

        # Chunk the text
        text_chunks = chunk_text_semantic(cleaned_text)

        for idx, chunk_text in enumerate(text_chunks):
            chunk = DocumentChunk(
                text=chunk_text,
                page_num=idx + 1,  # Simulated page number
                metadata={
                    "source": file_path,
                    "file_type": "docx"
                }
            )
            all_chunks.append(chunk)

        return all_chunks


class ExcelLoader(DocumentLoader):
    """Loader for Excel and CSV files"""

    SUPPORTED_EXTENSIONS = {'.xlsx', '.xls', '.csv'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in ExcelLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load Excel/CSV file and extract text chunks"""
        if pd is None:
            raise ImportError("pandas is required for Excel/CSV support. Install with: pip install pandas openpyxl")

        file_ext = Path(file_path).suffix.lower()
        all_chunks = []

        # Read file
        if file_ext == '.csv':
            try:
                # Try reading as CSV with different encodings
                try:
                    df = pd.read_csv(file_path, encoding='utf-8')
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding='latin1')
                sheets = {'Sheet1': df}
            except Exception as e:
                print(f"Error reading CSV {file_path}: {e}")
                return []
        else:
            # Read all sheets - let pandas auto-detect engine ('openpyxl' for xlsx, 'xlrd' for xls)
            # NOTE: explicit engine='openpyxl' forces XML parsing which fails on binary .xls or corrupted files
            try:
                sheets = pd.read_excel(file_path, sheet_name=None)
            except Exception as e:
                error_msg = str(e)
                print(f"Error reading Excel {file_path}: {error_msg}")
                
                # Check for common "No valid workbook part" (corrupted/encrypted xlsx)
                if "valid workbook part" in error_msg:
                   raise ValueError(f"Excel file is corrupted or encrypted. Please re-save it as a standard .xlsx file. (Error: {error_msg})")
                
                # Check for missing dependency for .xls
                if "xlrd" in error_msg or "optional dependency" in error_msg:
                     raise ImportError("Legacy .xls support requires 'xlrd'. Please install it: pip install xlrd>=2.0.1")
                
                raise e

        # Process each sheet
        for sheet_name, df in sheets.items():
            # Convert dataframe to readable text
            text_parts = []

            # Add column headers
            text_parts.append("Columns: " + ", ".join(df.columns.astype(str)))
            text_parts.append("")

            # Add rows (limit to prevent massive chunks)
            for idx, row in df.head(1000).iterrows():  # Limit to first 1000 rows
                row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                if row_text:
                    text_parts.append(row_text)

            combined_text = "\n".join(text_parts)
            cleaned_text = clean_text(combined_text)

            if not cleaned_text:
                continue

            # Chunk the text
            text_chunks = chunk_text_semantic(cleaned_text)

            for idx, chunk_text in enumerate(text_chunks):
                chunk = DocumentChunk(
                    text=chunk_text,
                    page_num=idx + 1,
                    metadata={
                        "source": file_path,
                        "sheet_name": sheet_name,
                        "file_type": file_ext[1:]  # Remove dot
                    }
                )
                all_chunks.append(chunk)

        return all_chunks


class HTMLLoader(DocumentLoader):
    """Loader for HTML and XML files"""

    SUPPORTED_EXTENSIONS = {'.html', '.htm', '.xml'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in HTMLLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load HTML/XML file and extract text chunks"""
        if BeautifulSoup is None:
            raise ImportError("beautifulsoup4 is required for HTML/XML support. Install with: pip install beautifulsoup4 lxml")

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # Parse HTML/XML
        soup = BeautifulSoup(content, 'lxml')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text
        text = soup.get_text()
        cleaned_text = clean_text(text)

        if not cleaned_text:
            return []

        # Chunk the text
        text_chunks = chunk_text_semantic(cleaned_text)
        all_chunks = []

        for idx, chunk_text in enumerate(text_chunks):
            chunk = DocumentChunk(
                text=chunk_text,
                page_num=idx + 1,
                metadata={
                    "source": file_path,
                    "file_type": Path(file_path).suffix[1:]
                }
            )
            all_chunks.append(chunk)

        return all_chunks


class TextLoader(DocumentLoader):
    """Loader for plain text, code, and markdown files"""

    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.markdown', '.rst',
        # Code files
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.c', '.cpp', '.cs',
        '.go', '.rs', '.php', '.rb', '.sql', '.sh', '.bash', '.json', '.yaml', '.yml',
        '.toml', '.ini', '.cfg', '.conf', '.env'
    }

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in TextLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load text-based file and extract chunks"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        except Exception as e:
            print(f"Error reading file {file_path}: {str(e)}")
            return []

        cleaned_text = clean_text(text)

        if not cleaned_text:
            return []

        # Chunk the text
        text_chunks = chunk_text_semantic(cleaned_text)
        all_chunks = []

        file_ext = Path(file_path).suffix

        for idx, chunk_text in enumerate(text_chunks):
            chunk = DocumentChunk(
                text=chunk_text,
                page_num=idx + 1,
                metadata={
                    "source": file_path,
                    "file_type": file_ext[1:] if file_ext else "txt"
                }
            )
            all_chunks.append(chunk)

        return all_chunks


class ImageLoader(DocumentLoader):
    """Loader for image files (PNG, JPG, JPEG, WEBP)"""

    SUPPORTED_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in ImageLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """Load image file and convert to base64"""
        if Image is None:
            raise ImportError("Pillow is required for image support. Install with: pip install Pillow")

        try:
            # Open and convert to RGB
            img = Image.open(file_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')

            # Resize if too large (max 1024px on any dimension) - safer for preview models
            max_dim = 1024
            if img.width > max_dim or img.height > max_dim:
                ratio = min(max_dim / img.width, max_dim / img.height)
                new_size = (int(img.width * ratio), int(img.height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)

            # Convert to base64
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=85)
            img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

            # Create chunk with image
            chunk = DocumentChunk(
                text="[Image File]",
                page_num=1,
                metadata={
                    "source": file_path,
                    "file_type": Path(file_path).suffix[1:],
                    "width": img.width,
                    "height": img.height
                },
                images=[img_base64]
            )

            return [chunk]

        except Exception as e:
            print(f"Error loading image {file_path}: {str(e)}")
            return []


class PPTLoader(DocumentLoader):
    """Loader for PowerPoint presentations (.pptx)"""

    SUPPORTED_EXTENSIONS = {'.pptx', '.ppt'}

    @staticmethod
    def supports(file_extension: str) -> bool:
        return file_extension.lower() in PPTLoader.SUPPORTED_EXTENSIONS

    def load(self, file_path: str) -> List[DocumentChunk]:
        """
        Load PowerPoint file and extract text chunks with embedded images.
        Aggregates text from all slides to prevent small-chunk skipping.
        """
        if Presentation is None:
            raise ImportError("python-pptx is required for PowerPoint support. Install with: pip install python-pptx")

        if file_path.lower().endswith('.ppt'):
             raise ValueError("Legacy .ppt files are not supported by python-pptx. Please convert to .pptx")

        prs = Presentation(file_path)

        full_presentation_text = ""
        all_images = []

        # Extract text and images from slides and aggregate
        for idx, slide in enumerate(prs.slides):
            slide_header = f"--- Slide {idx + 1} ---\n"
            slide_content = []

            # Extract from shapes (text boxes, images, etc)
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

                # Extract images
                if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        
                        # Process image with Pillow (Resize + Convert to JPEG)
                        if Image:
                            img = Image.open(io.BytesIO(image_blob))
                            if img.mode != 'RGB':
                                img = img.convert('RGB')
                            
                            # Limit size (same as ImageLoader)
                            max_dim = 1024
                            if img.width > max_dim or img.height > max_dim:
                                ratio = min(max_dim / img.width, max_dim / img.height)
                                new_size = (int(img.width * ratio), int(img.height * ratio))
                                img = img.resize(new_size, Image.Resampling.LANCZOS)
                                
                            buffer = io.BytesIO()
                            img.save(buffer, format='JPEG', quality=85)
                            img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
                        else:
                            # Fallback if Pillow is somehow missing (though we check imports)
                            # But this is risky as it might be PNG
                            img_base64 = base64.b64encode(image_blob).decode('utf-8')

                        all_images.append(img_base64)
                        slide_content.append(f"[Image {len(all_images)}]")
                    except Exception as e:
                        print(f"Warning: Could not extract image from slide {idx + 1}: {str(e)}")

            # Extract from notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                 notes = slide.notes_slide.notes_text_frame.text
                 if notes.strip():
                     slide_content.append(f"[Notes]: {notes}")

            if slide_content:
                full_presentation_text += slide_header + "\n".join(slide_content) + "\n\n"

        cleaned_text = clean_text(full_presentation_text)

        if not cleaned_text and not all_images:
            return []

        # If we only have images but no text, create a placeholder
        if not cleaned_text:
            cleaned_text = "[PowerPoint Presentation with Images]"

        # Chunk the entire presentation text semantically
        text_chunks = chunk_text_semantic(cleaned_text)
        all_chunks = []

        # Attach all images to the first chunk (since we're treating as single flow)
        for idx, chunk_text in enumerate(text_chunks):
            chunk = DocumentChunk(
                text=chunk_text,
                page_num=1,  # Treat as a single document flow
                metadata={
                    "source": file_path,
                    "file_type": "pptx",
                    "total_slides": len(prs.slides),
                    "image_count": len(all_images)
                },
                images=all_images if idx == 0 else None  # Attach all images to first chunk
            )
            all_chunks.append(chunk)

        # DEBUG: Print chunk count
        print(f"[PPTLoader] Extracted {len(all_chunks)} chunks from {len(prs.slides)} slides ({len(all_images)} images).")

        return all_chunks


# Factory function to get appropriate loader
def get_loader_for_file(file_path: str) -> DocumentLoader:
    """
    Get the appropriate loader for a file based on its extension.

    Args:
        file_path: Path to the file

    Returns:
        DocumentLoader instance

    Raises:
        ValueError: If file type is not supported
    """
    file_ext = Path(file_path).suffix.lower()

    # Try each loader
    loaders = [PDFLoader, WordLoader, ExcelLoader, HTMLLoader, ImageLoader, PPTLoader, TextLoader]

    for loader_class in loaders:
        if loader_class.supports(file_ext):
            return loader_class()

    raise ValueError(f"Unsupported file type: {file_ext}")


def process_document(file_path: str) -> List[DocumentChunk]:
    """
    Process any supported document type and return chunks.
    This is the main entry point that replaces process_pdf.

    Args:
        file_path: Path to the document file

    Returns:
        List of DocumentChunk objects
    """
    loader = get_loader_for_file(file_path)
    return loader.load(file_path)
